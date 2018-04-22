#!/usr/bin/env python3
import pptx
import googlesearch

# create a presentation
my_presentation = pptx.Presentation('/Users/megandu/Downloads/CH17_Lecture.pptx')
things_to_search = [] # create a list to store all of the titles

# iterate over the slides
for slide in my_presentation.slides:
    # titles should be stored in shapes.title and the string in .text
    title = slide.shapes.title.text

    # make sure it's not None
    if title:
      things_to_search.append(title) # if it exists, append it to our list
    else:
      print('No title found in slide')


# iterate over our list of titles to search
for thing in things_to_search:
    # make sure that we want to search for it
    response = input('Got title: "' + thing + '" do you want to search for results? ')
    if response == 'y' or response == 'Y':
        for url in googlesearch.search(thing + ' youtube', stop=5):
                print('\t' + url)
