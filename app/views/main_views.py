# Copyright 2014 SolidBuilds.com. All rights reserved
#
# Authors: Ling Thio <ling.thio@gmail.com>


from flask import Blueprint, redirect, render_template
from flask import request, url_for, send_from_directory
from flask_user import current_user, login_required, roles_required

from app import db

from app.models.user_models import UserProfileForm, User, RegisteredClass, AddClassForm, DeleteClassForm

from flask import send_from_directory
import pptx
import urllib.request
import urllib.parse
import re

from pptx import Presentation
import textract
import nltk
from nltk import sent_tokenize
import wikipedia
import requests
from bs4 import BeautifulSoup
from bs4.element import Comment

# Imports the Google Cloud client library
import argparse
from google.cloud import language
from google.cloud.language import enums
from google.cloud.language import types
from app.models.user_models import UserProfileForm, User, RegisteredClass, AddClassForm

main_blueprint = Blueprint('main', __name__, template_folder='templates')

# The Home page is accessible to anyone
@main_blueprint.route('/')
def home_page():
    return render_template('main/home_page.html')


# The User page is accessible to authenticated users (users that have logged in)
@main_blueprint.route('/dashboard', methods=['GET', 'POST'])
@login_required  # Limits access to authenticated users
def member_page():
    # Initialize form
    form = AddClassForm(request.form)

    # Process valid POST
    if request.method == 'POST' and form.validate():

        print(form.name.data, form)

        single_class = RegisteredClass(name=form.name.data, user_id=current_user.id)
        db.session.add(single_class)
        db.session.commit()

        # db.session.commit()
        return redirect(url_for('main.member_page'))

    classes = RegisteredClass.query.filter_by(user_id=current_user.id).all()
    return render_template('main/dashboard.html', classes = classes, form=form)

@main_blueprint.route('/upload')
@login_required  # Limits access to authenticated users
def upload_page():
    return render_template('main/upload.html')


# The Admin page is accessible to users with the 'admin' role
@main_blueprint.route('/admin')
@roles_required('admin')  # Limits access to users with the 'admin' role
def admin_page():
    return render_template('main/admin_page.html')

# @main_blueprint.route("/deleteClass", methods=["POST"])
# @login_required
# def delete_class():
#     form = DeleteClassForm(request.form)
#     id_delete = request.form.get("id")
#     print(id_delete)
#     # classes = RegisteredClass.query.filter_by(id_delete).first()
#     # db.session.delete(classes)
#     # db.session.commit()
#     return render_template('main/dashboard.html',
#                            form=form)

@main_blueprint.route('/classes/<int:id>', methods=['DELETE'])
def delete_entry(id):
    classes = RegisteredClass.query.filter_by(id).first()
    db.session.delete(classes)
    db.session.commit()
    return render_template('main/dashboard.html')


@main_blueprint.route('/main/profile', methods=['GET', 'POST'])
@login_required
def user_profile_page():
    # Initialize form
    form = UserProfileForm(request.form, obj=current_user)

    # Process valid POST
    if request.method == 'POST' and form.validate():
        # Copy form fields to user_profile fields
        form.populate_obj(current_user)

        # Save user_profile
        db.session.commit()

        # Redirect to home page
        return redirect(url_for('main.home_page'))

    # Process GET or invalid POST
    return render_template('main/user_profile_page.html',
                           form=form)

## you tube is banned for me right now so I can't test this
@main_blueprint.route('/youtube_info')
@login_required
def youtube_info():
    my_presentation = pptx.Presentation('/Users/jnai/Documents/hackathon_2018/app/static/files/test_2.pptx')
    phrases_to_search = [] # create a list to store all of the titles
    search_results = []

    for slide in my_presentation.slides:
        try:
            title = slide.shapes.title.text
        except:
            title = ''

        if title and title not in phrases_to_search:
          phrases_to_search.append(title)

    for phrase in phrases_to_search:
        try:
            query_string = urllib.parse.urlencode({"search_query" : phrase})
            html_content = urllib.request.urlopen("http://www.youtube.com/results?" + query_string)
            search_results = re.findall(r'href=\"\/watch\?v=(.{11})', html_content.read().decode())
            youtube_url = "http://www.youtube.com/watch?v=" + search_results[0]
            print(youtube_url)
        except:
            youtube_url = 'None'

        search_result = {
            'phrase': phrase,
            'youtube_link': youtube_url
        }
        print(search_result)

        search_results.append(search_result)
    return render_template('main/youtube_info.html')


@main_blueprint.route('/wiki_info')
@login_required
def wiki_info():
    # Getting visible text from wiki page
    def tag_visible(element):
        if element.parent.name in ['style', 'script', 'head', 'title', 'meta', '[document]']:
            return False
        if isinstance(element, Comment):
            return False
        return True

    def text_from_html(body):
        soup = BeautifulSoup(body, 'html.parser')
        texts = soup.findAll(text=True)
        visible_texts = filter(tag_visible, texts)
        return u" ".join(t.strip() for t in visible_texts)

    # Presi to text
    def remove_new_lines(all_text):
        filter = ['\\xe2', '\\x80', '\\x9', '\\x99t', '\\x93', '\\t', '\\xa6d']
        all_text = all_text.replace('\\n', ' ')
        all_text = all_text.replace('b\'', '') #weird
        for i in filter:
            all_text = all_text.replace(i, '')
        return all_text

    def pull_presi_text(filename):
        presi_in = textract.process(filename)
        all_text = remove_new_lines(str(presi_in))
        # sentences = sent_tokenize(all_text)
        # return sentences
        return all_text

    # Google NLP
    def print_entity_results(entity):
        # entity types from enums.Entity.Type
        entity_type = ('UNKNOWN', 'PERSON', 'LOCATION', 'ORGANIZATION',
                       'EVENT', 'WORK_OF_ART', 'CONSUMER_GOOD', 'OTHER')

        print('=' * 20)
        print(u'{:<16}: {}'.format('name', entity.name))
        print(u'{:<16}: {}'.format('type', entity_type[entity.type]))
        print(u'{:<16}: {}'.format('metadata', entity.metadata))
        print(u'{:<16}: {}'.format('salience', entity.salience))
        print(u'{:<16}: {}'.format('wikipedia_url',
              entity.metadata.get('wikipedia_url', '-')))

    def entities_text(text):
        """Detects entities in the text."""
        client = language.LanguageServiceClient()

        # if isinstance(text, six.binary_type):
        #     text = text.decode('utf-8')

        # Instantiates a plain text document.
        document = types.Document(
            content=text,
            type=enums.Document.Type.PLAIN_TEXT)

        # Detects entities in the document. You can also analyze HTML with:
        #   document.type == enums.Document.Type.HTML
        entities = client.analyze_entities(document).entities

        # for entity in entities:
        #     print_entity_results(entity)

        return entities

    # Getting wiki definitions
    def get_wiki_definitions(entities):
        phrases_and_defs = []

        for entity in entities:
            wiki_link = entity.metadata.get('wikipedia_url', '-')
            if wiki_link != '-':
                # print_entity_results(entity)
                try:
                    try:
                        wiki_search_term = wikipedia.search(entity.name)[0]
                        wiki_content = wikipedia.summary(wiki_search_term, sentences = 2)
                    except:
                        wiki_search_term = wikipedia.search(entity.name)[1]
                        wiki_content = wikipedia.summary(wiki_search_term, sentences = 2)

                    phrase_and_def = {
                        'phrase': entity.name,
                        'def': wiki_content
                    }
                    phrases_and_defs.append(phrase_and_def)
                    print(phrase_and_def)
                except:
                    pass

        return phrases_and_defs

    presi_text = pull_presi_text('/Users/jnai/Documents/hackathon_2018/app/static/files/test_2.pptx')
    entities = entities_text(presi_text)
    phrases_and_defs = get_wiki_definitions(entities)

    return render_template('main/wiki_info.html', wiki_info = phrases_and_defs)
